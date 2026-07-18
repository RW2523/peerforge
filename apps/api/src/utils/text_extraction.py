"""
Text extraction utilities for various file formats
Supports: PDF, DOCX, TXT, MD
"""

import io
import hashlib
from typing import Dict, List, Tuple
import filetype
from PyPDF2 import PdfReader
from docx import Document


class TextExtractor:
    """Extract text from uploaded files with provenance tracking"""
    
    # Allowed file types (MIME types)
    ALLOWED_TYPES = {
        'application/pdf': '.pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
        'text/plain': '.txt',
        'text/markdown': '.md',
    }

    # Audio types allowed only for meeting transcripts (transcribed before chunking)
    AUDIO_TYPES = {
        'audio/mpeg': '.mp3',
        'audio/mp4': '.m4a',
        'audio/x-m4a': '.m4a',
        'audio/wav': '.wav',
        'audio/x-wav': '.wav',
        'audio/webm': '.webm',
    }

    # Max file size: 50MB
    MAX_FILE_SIZE = 50 * 1024 * 1024

    @staticmethod
    def validate_file(file_data: bytes, filename: str, allow_audio: bool = False) -> Tuple[bool, str, str]:
        """
        Validate file type and size using magic bytes

        Args:
            file_data: File contents
            filename: Original filename
            allow_audio: When True, audio MIME types are also accepted (transcript uploads)

        Returns:
            (is_valid, mime_type, error_message)
        """
        # Check size
        if len(file_data) > TextExtractor.MAX_FILE_SIZE:
            return False, '', f'File size exceeds {TextExtractor.MAX_FILE_SIZE // (1024*1024)}MB limit'

        # Detect MIME type using magic bytes
        kind = filetype.guess(file_data)

        if kind is None:
            # Fall back to text/plain or text/markdown based on extension
            if filename.endswith('.txt'):
                mime_type = 'text/plain'
            elif filename.endswith('.md'):
                mime_type = 'text/markdown'
            else:
                return False, '', 'Unknown file type'
        else:
            mime_type = kind.mime

        allowed = dict(TextExtractor.ALLOWED_TYPES)
        if allow_audio:
            allowed.update(TextExtractor.AUDIO_TYPES)

        # Check if allowed
        if mime_type not in allowed:
            suffix = ' (audio only allowed for transcripts)' if mime_type in TextExtractor.AUDIO_TYPES else ''
            return False, mime_type, f'File type {mime_type} not allowed{suffix}'

        return True, mime_type, ''
    
    @staticmethod
    def extract_from_pdf(file_data: bytes) -> Dict:
        """
        Extract text from PDF with page-level provenance
        
        Returns:
            {
                'text': str,  # Full text
                'pages': List[Dict],  # Per-page data
                'page_count': int,
                'word_count': int,
                'is_scanned': bool,  # True if no text layer detected
                'sha256': str
            }
        """
        try:
            pdf_reader = PdfReader(io.BytesIO(file_data))
            pages = []
            full_text = []
            total_words = 0
            
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                text = page.extract_text()
                word_count = len(text.split())
                total_words += word_count
                
                pages.append({
                    'page_num': page_num,
                    'text': text,
                    'word_count': word_count,
                    'char_start': len(''.join(full_text)),
                    'char_end': len(''.join(full_text)) + len(text)
                })
                full_text.append(text)
            
            combined_text = '\n\n'.join(full_text)
            
            # Detect scanned PDF (very low word count suggests no text layer)
            is_scanned = total_words < 50 and len(pdf_reader.pages) > 1
            
            return {
                'text': combined_text,
                'pages': pages,
                'page_count': len(pdf_reader.pages),
                'word_count': total_words,
                'is_scanned': is_scanned,
                'sha256': hashlib.sha256(file_data).hexdigest(),
                'extraction_method': 'pypdf2'
            }
        
        except Exception as e:
            raise Exception(f'PDF extraction failed: {str(e)}')
    
    @staticmethod
    def extract_from_docx(file_data: bytes) -> Dict:
        """
        Extract text from DOCX with paragraph-level provenance
        
        Returns:
            {
                'text': str,
                'paragraphs': List[Dict],
                'paragraph_count': int,
                'word_count': int,
                'sha256': str
            }
        """
        try:
            doc = Document(io.BytesIO(file_data))
            paragraphs = []
            full_text = []
            total_words = 0
            
            for para_num, para in enumerate(doc.paragraphs, start=1):
                text = para.text.strip()
                if not text:
                    continue
                
                word_count = len(text.split())
                total_words += word_count
                
                paragraphs.append({
                    'paragraph_num': para_num,
                    'text': text,
                    'word_count': word_count,
                    'char_start': len('\n\n'.join(full_text)),
                    'char_end': len('\n\n'.join(full_text)) + len(text)
                })
                full_text.append(text)
            
            combined_text = '\n\n'.join(full_text)
            
            return {
                'text': combined_text,
                'paragraphs': paragraphs,
                'paragraph_count': len(paragraphs),
                'word_count': total_words,
                'sha256': hashlib.sha256(file_data).hexdigest(),
                'extraction_method': 'python-docx'
            }
        
        except Exception as e:
            raise Exception(f'DOCX extraction failed: {str(e)}')
    
    @staticmethod
    def extract_from_text(file_data: bytes) -> Dict:
        """
        Extract text from plain text or markdown files
        
        Returns:
            {
                'text': str,
                'word_count': int,
                'line_count': int,
                'sha256': str
            }
        """
        try:
            text = file_data.decode('utf-8')
            lines = text.split('\n')
            words = text.split()
            
            return {
                'text': text,
                'word_count': len(words),
                'line_count': len(lines),
                'sha256': hashlib.sha256(file_data).hexdigest(),
                'extraction_method': 'plain_text'
            }
        
        except UnicodeDecodeError:
            raise Exception('Text file encoding not supported (must be UTF-8)')
        except Exception as e:
            raise Exception(f'Text extraction failed: {str(e)}')
    
    @staticmethod
    def extract(file_data: bytes, mime_type: str) -> Dict:
        """
        Extract text based on MIME type
        
        Args:
            file_data: File contents
            mime_type: MIME type
        
        Returns:
            Extraction result dictionary
        
        Raises:
            Exception if extraction fails
        """
        if mime_type == 'application/pdf':
            return TextExtractor.extract_from_pdf(file_data)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return TextExtractor.extract_from_docx(file_data)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return TextExtractor.extract_from_pptx(file_data)
        elif mime_type in ['text/plain', 'text/markdown']:
            return TextExtractor.extract_from_text(file_data)
        else:
            raise Exception(f'Unsupported MIME type: {mime_type}')

    @staticmethod
    def extract_from_pptx(file_data: bytes) -> Dict:
        """
        Extract text from a PowerPoint deck with per-slide provenance.

        Emits per-slide entries in the same shape as PDF `pages` (page_num =
        slide number) so downstream chunking attaches a slide number to every
        chunk without any changes — provenance badges, Glass-Box links and
        certificates all key off the same field. Speaker notes are included
        inline, labelled, because they carry the presenter's actual narrative.

        Returns:
            {
                'text': str,
                'pages': List[{'page_num', 'text', 'char_start', 'char_end', 'word_count'}],
                'slides': List[{'slide_num', 'title', 'body_words', 'notes_words', 'bullet_count'}],
                'page_count': int, 'word_count': int, 'is_scanned': False, 'sha256': str,
                'extraction_method': 'python-pptx'
            }
        """
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_data))
        pages: List[Dict] = []
        slides_meta: List[Dict] = []
        parts: List[str] = []
        offset = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            bullet_count = 0
            title_shape = slide.shapes.title
            title = (title_shape.text or '').strip() if title_shape is not None else ''
            for shape in slide.shapes:
                if shape is title_shape or not getattr(shape, 'has_text_frame', False):
                    continue
                for para in shape.text_frame.paragraphs:
                    line = ''.join(run.text for run in para.runs).strip()
                    if not line:
                        continue
                    texts.append(line)
                    bullet_count += 1
            if not title and texts:
                title = texts[0]

            notes = ''
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = (slide.notes_slide.notes_text_frame.text or '').strip()

            body = '\n'.join(texts)
            slide_text = f"[Slide {slide_num}] {title}\n{body}"
            if notes:
                slide_text += f"\n[Speaker notes] {notes}"

            char_start = offset
            char_end = offset + len(slide_text)
            pages.append({
                'page_num': slide_num,
                'text': slide_text,
                'char_start': char_start,
                'char_end': char_end,
                'word_count': len(slide_text.split()),
            })
            slides_meta.append({
                'slide_num': slide_num,
                'title': title[:120],
                'body_words': len(body.split()),
                'notes_words': len(notes.split()),
                'bullet_count': bullet_count,
            })
            parts.append(slide_text)
            offset = char_end + 2  # account for the joining "\n\n"

        full_text = '\n\n'.join(parts)
        return {
            'text': full_text,
            'pages': pages,
            'slides': slides_meta,
            'page_count': len(pages),
            'word_count': len(full_text.split()),
            'is_scanned': False,
            'sha256': hashlib.sha256(file_data).hexdigest(),
            'extraction_method': 'python-pptx',
        }
